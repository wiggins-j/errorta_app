"""PPTX extraction via python-pptx. One chunk per slide."""
from __future__ import annotations

from pathlib import Path

from . import Chunk, ExtractError


def extract(path: Path) -> list[Chunk]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise ExtractError(f"python-pptx not available: {e}") from e

    try:
        pres = Presentation(str(path))
    except Exception as e:
        raise ExtractError(f"could not open PPTX: {e}") from e

    chunks: list[Chunk] = []
    for i, slide in enumerate(pres.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
        # Speaker notes
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    parts.append(f"[notes] {note}")
        except Exception:
            pass
        text = "\n".join(parts).strip()
        if not text:
            continue
        chunks.append(
            {"text": text, "meta": {"slide_number": i + 1, "source_type": "pptx"}}
        )
    if not chunks:
        raise ExtractError("PPTX has no text content")
    return chunks
